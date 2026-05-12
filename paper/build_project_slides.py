#!/usr/bin/env python3
"""Generate GNN privacy project deck: ≤20 slides, clean designed layout."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "GNN_Privacy_Project_Overview.pptx"

# Design system
SLATE_DEEP = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x14, 0xB8, 0xA6)
BG = RGBColor(0xF8, 0xFA, 0xFC)
BODY = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
FOOTER_TEXT = "Graph neural networks & privacy — study overview"


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_slide_background(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_content_slide(prs, title: str, bullets: list[str], *, body_pt=19):
    """Dark header bar, teal accent rule, light body, footer."""
    s = blank_slide(prs)
    w, h = prs.slide_width, prs.slide_height
    set_slide_background(s, BG)

    bar_h = Inches(1.14)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = SLATE_DEEP
    bar.line.fill.background()

    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, bar_h, w, Inches(0.07))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()

    tbox = s.shapes.add_textbox(Inches(0.65), Inches(0.36), w - Inches(1.3), Inches(0.75))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE

    body_top = Inches(1.45)
    body_h = Inches(5.35)
    bbox = s.shapes.add_textbox(Inches(0.65), body_top, w - Inches(1.3), body_h)
    btf = bbox.text_frame
    btf.word_wrap = True
    btf.auto_size = MSO_AUTO_SIZE.NONE
    for i, line in enumerate(bullets):
        bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        bp.text = line
        bp.level = 0
        bp.font.size = Pt(body_pt)
        bp.font.color.rgb = BODY
        bp.font.name = FONT_BODY
        bp.space_after = Pt(14)
        bp.line_spacing = 1.15

    fbox = s.shapes.add_textbox(Inches(0.65), h - Inches(0.52), w - Inches(1.3), Inches(0.35))
    ftf = fbox.text_frame
    fp = ftf.paragraphs[0]
    fp.text = FOOTER_TEXT
    fp.font.size = Pt(11)
    fp.font.color.rgb = MUTED
    fp.font.name = FONT_BODY


def add_title_slide(prs):
    s = blank_slide(prs)
    w, h = prs.slide_width, prs.slide_height
    set_slide_background(s, BG)

    hero = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.05), w, Inches(3.15))
    hero.fill.solid()
    hero.fill.fore_color.rgb = SLATE_DEEP
    hero.line.fill.background()

    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.05), Inches(0.18), Inches(3.15))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    tb = s.shapes.add_textbox(Inches(0.95), Inches(2.35), w - Inches(1.5), Inches(2.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Privacy Leakage & Lightweight Defenses"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE
    p2 = tf.add_paragraph()
    p2.text = "Graph Neural Networks on Sensitive Relational Data"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xCC, 0xFB, 0xF1)
    p2.font.name = FONT_BODY
    p2.space_before = Pt(8)

    sub = s.shapes.add_textbox(Inches(0.65), Inches(5.55), w - Inches(1.3), Inches(0.9))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sp.text = "An empirical study on membership privacy and simple defenses"
    sp.font.size = Pt(17)
    sp.font.color.rgb = BODY
    sp.font.name = FONT_BODY
    sp.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_content_slide(
        prs,
        "What this study is about",
        [
            "Graph neural networks learn from networks of connected points (people, papers, devices). The question: if you only see the model’s predicted probabilities for each point, can you guess whether that point was used in training?",
            "That guess is called a membership inference attack. If attackers do better than random guessing, the model is leaking information about who was in the training set.",
            "We built artificial networks where we can dial two knobs on purpose: how much “similar points link together,” and how crowded the graph is with links—so we are not stuck with whatever a single real dataset happens to look like.",
            "We also tried several cheap, practical training tweaks people might use as privacy fixes, and checked whether they actually helped or accidentally made things worse.",
        ],
    )

    add_content_slide(
        prs,
        "Why we care",
        [
            "Many real graphs describe sensitive relationships. If a deployed model reveals training membership, that can harm people or organizations even when raw data stay private.",
            "In practice, services often return confidence-style scores, not just a single label. Those scores are exactly what these attacks use—no need to see inside the model.",
            "Because the model mixes information across neighbors, risk might depend on the shape of the graph—not only on how accurate the model looks on a test set.",
        ],
    )

    add_content_slide(
        prs,
        "Questions we asked",
        [
            "Do graph models leak more than strong models that ignore the graph but use the same inputs?",
            "When neighbors are “alike” vs “mixed,” and when the graph is sparse vs dense, how does attack success change?",
            "Which simple defenses really cut leakage without ruining accuracy?",
            "Where do we get the best balance of usefulness vs privacy? (We summarize attack success with a score where 0.5 means random guessing and higher means easier to spot members.)",
        ],
    )

    add_content_slide(
        prs,
        "How the study was organized",
        [
            "Everything is driven from a single settings file, so runs stay organized and repeatable—no one-off tweaks hidden in code.",
            "The main experiment grid covers six graph types, four model types, several defense options, and five random restarts—hundreds of full train-and-test cycles in total.",
            "Each run saves results to tables; we average across restarts and use simple statistical checks when comparing a defense to “no defense.”",
        ],
    )

    add_content_slide(
        prs,
        "The graphs we used",
        [
            "We mainly used synthetic graphs: 400 points, 50 input numbers per point, 5 categories. Half the points train the model; half are held out—reshuffled on each restart.",
            "“Homophily” here means: do links mostly connect points in the same category? We study both “mostly yes” and “mostly no,” crossed with sparse, medium, or busy link patterns.",
            "Small citation-style datasets and larger graphs can also be plugged into the same code; the written results focus on the synthetic grid so the two knobs are clean.",
        ],
    )

    add_content_slide(
        prs,
        "Models we compared",
        [
            "Two baselines only look at each point’s own features—classic logistic regression and a small neural network. They tell us how much signal comes from features alone.",
            "A graph convolutional network (GCN) passes information along links each layer—good when neighbors are helpful, trickier when links are noisy.",
            "GraphSAGE also uses neighbors but in a different way; in our study it often matched GCN on accuracy while behaving much more like the baselines from a privacy angle.",
        ],
    )

    add_content_slide(
        prs,
        "How attacks and “calibration” were measured",
        [
            "We assume an outsider sees the full list of category probabilities the model outputs for a point, and knows the true category being tested—typical for node-level membership games.",
            "Three attack styles: a small classifier built from simple summaries of those probabilities; a one-number rule based on confidence in the true category; and a “shadow model” trained separately to mimic the same setup.",
            "We also measure calibration: when the model says “about 80% sure,” is it right about 80% of the time on held-out data? Miscalibration often shows up alongside easier membership guesses.",
        ],
    )

    add_content_slide(
        prs,
        "Defenses we tried",
        [
            "Baseline training, then: randomly drop links during training; soften hard labels; stop early if training loss plateaus; at prediction time keep only the top two category scores; or permanently delete a fraction of links before training.",
            "Rough settings: drop about 30% of edges per training step when that defense is on; soften labels with a 10% mix toward uniform; keep the top 2 scores when masking; remove about 20% of edges once for the sparsification run.",
            "Graph models use these toggles; the feature-only baselines stay vanilla so we can separate “graph tricks” from “feature-only” behavior.",
        ],
        body_pt=17,
    )

    add_content_slide(
        prs,
        "Training and checking results",
        [
            "Standard optimizer settings, a modest epoch budget, and the same recipe across runs so differences come from the graph type, model, or defense—not from hand-tuning.",
            "Multiple random restarts mean we are not trusting one lucky split; a separate synthetic draw powers the shadow-model story so it is not copied from the main graph.",
            "We report averages and spread across restarts; significance notes are helpful but based on only five paired runs, so they are indicative, not the last word.",
        ],
    )

    add_content_slide(
        prs,
        "Takeaway: risk depends on the setting",
        [
            "There is no simple rule that “graph models always leak more.” When neighbors strongly match their categories, every model we tried stayed close to random guessing for membership.",
            "When neighbors are mixed and the graph is thinner, one architecture (GCN) became the easiest to attack; the other graph model and the non-graph models stayed nearer to coin-flip guessing.",
            "So privacy is about the combination of graph structure and model design—not the word “graph” by itself.",
        ],
    )

    add_content_slide(
        prs,
        "When the graph convolution model struggled most",
        [
            "The clearest hotspot was GCN on mixed-category, sparse graphs: the membership score rose to about 0.59, while others in the same conditions stayed near 0.52.",
            "Plain-language read: the model’s confidence patterns on training points looked more “special” than on held-out points, so a simple attacker could sort them better than chance.",
            "When links mostly agree with categories, that gap shrinks and all models look safer under this style of attack.",
        ],
    )

    add_content_slide(
        prs,
        "What happened to the defenses",
        [
            "Softening labels often backfired: attack scores went up and calibration got worse in several cases—so it should not be assumed to be a privacy fix without testing.",
            "Deleting a chunk of edges once before training was the one tweak with a clear, repeatable win for GCN in a “mixed categories, medium busyness” graph—attack score dropped in a statistically noticeable way.",
            "Random edge dropping each step sometimes helped a little; early stopping and hiding low scores did not show a steady privacy benefit across the board.",
        ],
    )

    add_content_slide(
        prs,
        "Accuracy is not the whole story",
        [
            "Two models can have almost the same accuracy but very different membership risk—what matters is how their probability outputs differ between training and held-out points.",
            "When we trained a separate shadow model and repeated the attack, the story matched the simpler attacks—so the conclusion is not an artifact of one method.",
            "Watching calibration alongside accuracy is useful: overconfident or unevenly confident models tend to be easier to attack in our results.",
        ],
    )

    add_content_slide(
        prs,
        "Practical guidance",
        [
            "Know your graph: if categories are mixed and the network is sparse, treat privacy as a first-class metric and actually run an attack-based check, especially for GCN-style models.",
            "Do not treat “softer labels” as a privacy tool by default—here it often made leakage worse.",
            "If you need a cheap knob to try before heavy formal privacy methods, thinning the graph before training is the one that showed a solid benefit in our hardest GCN case.",
        ],
    )

    add_content_slide(
        prs,
        "Caveats and next steps",
        [
            "Synthetic graphs are clean for science but simpler than messy real networks; follow-up work should stress-test the same questions on real data.",
            "We studied probability-based attacks; stronger or different adversaries could change the picture.",
            "The software side can add more restarts, extra attack types, and larger graphs; those extensions are meant for future runs, not the core story of this deck.",
        ],
    )

    add_content_slide(
        prs,
        "Software and reproducibility",
        [
            "A single configuration file lists datasets, models, defenses, random seeds, and which attacks to run—change the file, rerun, get new tables.",
            "The pipeline loads data, trains the target model, builds attack scores, and writes spreadsheets and plots so figures trace back to numbers.",
            "The goal is that another researcher can rerun or extend the grid without reverse-engineering one-off scripts.",
        ],
    )

    add_content_slide(
        prs,
        "Bottom line",
        [
            "Membership risk for node-level graph models is not universal: one architecture stood out when the graph was unfriendly, while another graph model often looked like a plain feature model from the attacker’s perspective.",
            "Simple defenses are not interchangeable; one helped in a targeted way, while another popular trick often hurt.",
            "The project pairs a clear experimental story with a repeatable workflow for trying new graphs, models, and defenses.",
        ],
    )

    add_content_slide(
        prs,
        "Write-ups and artifacts",
        [
            "A formal paper-style document spells out notation, the full grid, and all figures for readers who want every detail.",
            "A separate long-form summary explains the same results in everyday language for teaching or outreach.",
            "Together with the configuration-driven runs, these artifacts document what was done and how to reproduce or extend it.",
        ],
    )

    s = blank_slide(prs)
    w, h = prs.slide_width, prs.slide_height
    set_slide_background(s, BG)
    bar_h = Inches(1.14)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = SLATE_DEEP
    bar.line.fill.background()
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, bar_h, w, Inches(0.07))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    tbox = s.shapes.add_textbox(Inches(0.65), Inches(0.4), w - Inches(1.3), Inches(0.75))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = "Thank you"
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.font.name = FONT_TITLE
    cbox = s.shapes.add_textbox(Inches(0.65), Inches(2.5), w - Inches(1.3), Inches(3.5))
    ctf = cbox.text_frame
    for i, line in enumerate(["Questions?", "Thank you for listening."]):
        cp = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        cp.text = line
        cp.font.size = Pt(22)
        cp.font.color.rgb = BODY
        cp.font.name = FONT_BODY
        cp.alignment = PP_ALIGN.CENTER
        cp.space_after = Pt(16)
    fbox = s.shapes.add_textbox(Inches(0.65), h - Inches(0.52), w - Inches(1.3), Inches(0.35))
    fp = fbox.text_frame.paragraphs[0]
    fp.text = FOOTER_TEXT
    fp.font.size = Pt(11)
    fp.font.color.rgb = MUTED
    fp.font.name = FONT_BODY

    assert len(prs.slides) == 20, len(prs.slides)
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
