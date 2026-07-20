#!/usr/bin/env python3
"""Build professor overview PPTX for GNN membership-inference project (with speaker notes)."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

FONT = "Calibri"

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "figures"
OUT = Path(__file__).resolve().parent / "GNN_Privacy_Professor_Overview_FINAL.pptx"

C_TITLE = RGBColor(15, 23, 42)
C_BODY = RGBColor(51, 65, 85)
C_ACCENT = RGBColor(37, 99, 235)
C_MUTED = RGBColor(100, 116, 139)
C_BAR = RGBColor(37, 99, 235)
C_TAKEAWAY_BG = RGBColor(219, 234, 254)
C_TAKEAWAY_TEXT = RGBColor(30, 58, 138)


def _set_slide_size(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def _set_notes(slide, text: str) -> None:
    """Speaker notes visible in PowerPoint Notes pane / Presenter view."""
    t = text.strip()
    if not t:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = t


def _style_title(shape) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.font.name = FONT
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = C_TITLE
        p.alignment = PP_ALIGN.LEFT


def _add_accent_bar(slide, top=Inches(0), height=Inches(7.5)) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        top,
        Inches(0.22),
        height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_BAR
    bar.line.fill.background()


def title_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    opener_takeaway: str,
    notes: str,
) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _add_accent_bar(slide)
    box = slide.shapes.add_textbox(Inches(0.55), Inches(1.55), Inches(12.2), Inches(2.2))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.name = FONT
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C_TITLE
    sub = slide.shapes.add_textbox(Inches(0.55), Inches(3.75), Inches(12.2), Inches(1.35))
    stf = sub.text_frame
    stf.text = subtitle
    sp = stf.paragraphs[0]
    sp.font.name = FONT
    sp.font.size = Pt(17)
    sp.font.color.rgb = C_BODY
    tk = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(5.25),
        Inches(12.35),
        Inches(1.05),
    )
    try:
        tk.adjustments[0] = 0.08
    except (AttributeError, IndexError, TypeError):
        pass
    tk.fill.solid()
    tk.fill.fore_color.rgb = C_TAKEAWAY_BG
    tk.line.color.rgb = C_ACCENT
    tk.line.width = Pt(1)
    tkf = tk.text_frame
    tkf.clear()
    tkf.margin_left = Inches(0.2)
    tkf.margin_right = Inches(0.2)
    tkf.margin_top = Inches(0.08)
    tkf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tp = tkf.paragraphs[0]
    tp.text = opener_takeaway
    tp.font.name = FONT
    tp.font.size = Pt(20)
    tp.font.bold = True
    tp.font.color.rgb = C_TAKEAWAY_TEXT
    tp.alignment = PP_ALIGN.LEFT
    foot = slide.shapes.add_textbox(Inches(0.55), Inches(6.85), Inches(12), Inches(0.4))
    ff = foot.text_frame
    ff.text = "YAML-driven benchmark · Cora, Citeseer, controlled synthetics"
    ff.paragraphs[0].font.size = Pt(12)
    ff.paragraphs[0].font.color.rgb = C_MUTED
    ff.paragraphs[0].font.name = FONT
    _set_notes(slide, notes)


def takeaway_slide(
    prs: Presentation,
    headline: str,
    notes: str,
    supporting: str | None = None,
) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _add_accent_bar(slide)
    tag = slide.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(11.5), Inches(0.45))
    ttf = tag.text_frame
    ttf.text = "Takeaway"
    ttf.paragraphs[0].font.name = FONT
    ttf.paragraphs[0].font.size = Pt(14)
    ttf.paragraphs[0].font.bold = True
    ttf.paragraphs[0].font.color.rgb = C_ACCENT
    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.85), Inches(12.1), Inches(3.2))
    btf = body.text_frame
    btf.word_wrap = True
    btf.text = headline
    bp = btf.paragraphs[0]
    bp.font.name = FONT
    bp.font.size = Pt(34)
    bp.font.bold = True
    bp.font.color.rgb = C_TITLE
    bp.alignment = PP_ALIGN.LEFT
    bp.line_spacing = 1.15
    if supporting:
        sp = btf.add_paragraph()
        sp.text = supporting
        sp.font.name = FONT
        sp.font.size = Pt(20)
        sp.font.color.rgb = C_BODY
        sp.font.bold = False
        sp.space_before = Pt(16)
        sp.line_spacing = 1.2
    _set_notes(slide, notes)


def section_slide(prs: Presentation, label: str, title: str, notes: str) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _add_accent_bar(slide)
    lab = slide.shapes.add_textbox(Inches(0.55), Inches(2.5), Inches(12), Inches(0.5))
    ltf = lab.text_frame
    ltf.text = label.upper()
    ltf.paragraphs[0].font.size = Pt(14)
    ltf.paragraphs[0].font.color.rgb = C_ACCENT
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.name = FONT
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(3.05), Inches(12), Inches(1.5))
    tf = tb.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = C_TITLE
    tf.paragraphs[0].font.name = FONT
    _set_notes(slide, notes)


def bullet_slide(prs: Presentation, title: str, bullets: list[str], notes: str) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _style_title(slide.shapes.title)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    tf.margin_bottom = Inches(0.12)
    tf.margin_top = Inches(0.08)
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(21)
        p.font.color.rgb = C_BODY
        p.space_after = Pt(11)
    _set_notes(slide, notes)


def _add_takeaway_banner(slide, text: str, top: float, height: float = 0.82) -> None:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.45),
        Inches(top),
        Inches(12.4),
        Inches(height),
    )
    try:
        box.adjustments[0] = 0.06
    except (AttributeError, IndexError):
        pass
    box.fill.solid()
    box.fill.fore_color.rgb = C_TAKEAWAY_BG
    box.line.color.rgb = C_ACCENT
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.06)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = f"Takeaway: {text}"
    p.font.name = FONT
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = C_TAKEAWAY_TEXT
    p.alignment = PP_ALIGN.LEFT


def figure_slide(
    prs: Presentation,
    title: str,
    rel_path: str,
    takeaway: str,
    notes: str,
    caption: str | None = None,
) -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _style_title(slide.shapes.title)
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(26)

    _add_takeaway_banner(slide, takeaway, top=1.02, height=0.88)

    path = FIG / rel_path
    pic_left = Inches(0.55)
    pic_top = Inches(2.05)
    pic_w = Inches(12.2)
    if not path.exists():
        tb = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1))
        tb.text_frame.text = f"(Missing figure: {rel_path})"
        _set_notes(slide, notes + f"\n\n[Build note: missing file {path}]")
        return

    slide.shapes.add_picture(str(path), pic_left, pic_top, width=pic_w)

    if caption:
        cap = slide.shapes.add_textbox(Inches(0.55), Inches(6.55), Inches(12.25), Inches(0.85))
        ctf = cap.text_frame
        ctf.word_wrap = True
        ctf.text = caption
        ctf.paragraphs[0].font.size = Pt(12)
        ctf.paragraphs[0].font.color.rgb = C_MUTED
        ctf.paragraphs[0].font.name = FONT

    _set_notes(slide, notes)


def two_caption_slide(
    prs: Presentation,
    title: str,
    takeaway: str,
    left: tuple[str, str],
    right: tuple[str, str],
    notes: str,
) -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _style_title(slide.shapes.title)
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)

    _add_takeaway_banner(slide, takeaway, top=0.98, height=0.82)

    w = Inches(6.05)
    top = Inches(1.95)
    for i, (fname, cap) in enumerate((left, right)):
        path = FIG / fname
        left_pos = Inches(0.48 + i * 6.38)
        if path.exists():
            slide.shapes.add_picture(str(path), left_pos, top, width=w)
        cbox = slide.shapes.add_textbox(left_pos, Inches(5.95), w, Inches(1.15))
        ctf = cbox.text_frame
        ctf.text = cap
        ctf.paragraphs[0].font.size = Pt(11)
        ctf.paragraphs[0].font.color.rgb = C_MUTED
        ctf.paragraphs[0].font.name = FONT

    _set_notes(slide, notes)


def main() -> int:
    prs = Presentation()
    _set_slide_size(prs)

    title_slide(
        prs,
        "Structural drivers of membership inference in GNNs",
        "Can an adversary tell—using only output probabilities—whether a node’s label was in the training set?",
        "Privacy here is not “accuracy.” It is whether train vs held-out nodes look different in probability space.",
        notes=(
            "Open by grounding the project in a question professors care about: not graph learning accuracy alone, "
            "but whether deployed APIs leak supervision information. Clarify membership here means the node’s label "
            "was included in the training loss mask on a fixed transductive graph—the graph and features are still "
            "shared. Mention you sweep synthetic regimes plus Cora/Citeseer, and that all experiments are reproducible "
            "from YAML and CSV outputs."
        ),
    )

    bullet_slide(
        prs,
        "What are GNNs—and what privacy problem do we study?",
        [
            "GNNs learn a representation for each node by repeatedly mixing its features with its neighbors’ (message passing over edges).",
            "They fit relational data: citations, social networks, fraud graphs, and other settings where connections carry signal.",
            "Privacy concern: a client who only sees output probabilities may still infer whether a node’s label was in the training loss—membership inference—not the same as guessing labels or stealing the graph.",
            "This project measures how often that succeeds and whether simple defenses help, across graph structures and model types.",
        ],
        notes=(
            "Assume some listeners have never seen a GNN. Give a one-sentence intuition: each layer aggregates neighbors, "
            "so a node’s prediction depends on a multi-hop neighborhood. Then pivot to privacy: many deployments expose "
            "softmax scores via an API. Membership inference asks whether those scores look systematically different for "
            "nodes whose labels were supervised vs held out on the same graph. Emphasize transductive node classification: "
            "the graph is shared; only the training mask changes. That coupling is what makes graph membership non-trivial "
            "and structure-dependent."
        ),
    )

    takeaway_slide(
        prs,
        "We assume a realistic graph API: the adversary only sees softmax vectors (plus the node’s true label, as in standard membership games).",
        notes=(
            "Now that GNNs and membership inference are defined, tighten the adversary model. Emphasize that you intentionally avoid gradients, embeddings, "
            "or white-box access—this matches many production settings where clients only receive probabilities. The true "
            "label assumption follows common membership-inference formulations; if asked, note alternatives exist but "
            "this is the standard comparison point in the literature."
        ),
        supporting="If your product exposes probabilities, treat membership inference as part of the threat model.",
    )

    bullet_slide(
        prs,
        "What makes graphs different?",
        [
            "Neighbors are not independent—each node’s prediction depends on its neighborhood.",
            "We turn two knobs on purpose: homophily (same-label edges) and density (how many edges).",
            "GCN and GraphSAGE aggregate neighbors differently, so the same graph can yield different privacy risk.",
        ],
        notes=(
            "Contrast IID tabular learning: message passing couples nodes, so train and test nodes are not exchangeable "
            "even when only the mask changes. Homophily controls how often neighbors agree in label space; density controls "
            "how much mixing happens per hop. Preview that architecture changes the inductive bias—GCN behaves like a "
            "strong low-pass filter while GraphSAGE mean pooling responds differently under heterophily."
        ),
    )

    takeaway_slide(
        prs,
        "There is no one answer to “are GNNs private?” The answer depends on the graph regime, the model, and what you measure.",
        notes=(
            "This is the thesis slide before RQs. Push back on binary narratives like “GNNs always leak more.” Your "
            "results show the interaction of structure, model, and features. Invite the listener to think in terms of "
            "conditional statements: under what graph statistics and which architecture does the posterior gap appear?"
        ),
    )

    bullet_slide(
        prs,
        "What we ask (four questions)",
        [
            "Do GNNs always leak more than feature-only models?",
            "How do homophily and density change attack success?",
            "Do cheap defenses reduce leakage without ruining accuracy?",
            "Where do utility and privacy trade off?",
        ],
        notes=(
            "Walk through RQ1–RQ4 quickly. RQ1 compares hypothesis classes with similar feature dimensionality. RQ2 "
            "isolates generative factors you normally cannot separate in one real graph. RQ3 targets deployable defenses "
            "rather than DP guarantees. RQ4 is about whether you can be accurate yet still expose different membership "
            "signals—this becomes important for GraphSAGE vs GCN later."
        ),
    )

    bullet_slide(
        prs,
        "What we ran",
        [
            "Models: GCN & GraphSAGE vs LogReg & MLP (baselines ignore the adjacency).",
            "Synthetic graphs: 400 nodes—low/high homophily × sparse/medium/dense edges.",
            "Real benchmarks: Cora and Citeseer with the same attacks and defenses.",
            "Everything is YAML-driven with 5 seeds so results are reproducible.",
        ],
        notes=(
            "Give the experimental footprint: six synthetic regimes, two citation graphs, four model families, three "
            "attack families, six defense settings for GNNs (including none), five seeds. Mention experiment_config_paper.yaml "
            "as the paper grid. If asked about compute, note CPU sweeps are feasible at this scale."
        ),
    )

    takeaway_slide(
        prs,
        "Attacks are intentionally simple: four numbers derived from each node’s probability vector—enough to train a small classifier that guesses “member or not.”",
        notes=(
            "Motivate parsimony: you want leakage attributable to graph effects, not to an exotic attacker. The four φ "
            "features are standard interpretable signals (confidence, entropy, margin). Mention AUROC as the headline "
            "metric: 0.5 means the attacker cannot rank members above non-members better than chance."
        ),
        supporting="That keeps the story about the graph and the model—not about a fancy attacker.",
    )

    bullet_slide(
        prs,
        "How we attack (plain language)",
        [
            "Build 4 scores from the output: peak confidence, true-class probability, entropy, and a margin-style term.",
            "Main attack: logistic regression on those 4 scores (train/eval split inside the attack data).",
            "Also: a one-number threshold check, and a shadow-model transfer test.",
        ],
        notes=(
            "Explain the confidence attacker trains on half of the attack tuples and evaluates on the other half to "
            "avoid trivial overfitting of the meta-classifier. Shadow models use independent seeds or resampled splits "
            "for synthetic/citation settings to probe transfer. Threshold attack is a sanity baseline using only φ2."
        ),
    )

    bullet_slide(
        prs,
        "Defenses we tried (all lightweight)",
        [
            "DropEdge — randomly drop edges during training only.",
            "Label smoothing — soften one-hot targets.",
            "Early stopping — stop when training loss plateaus.",
            "Confidence masking — only show top-2 class probabilities at inference.",
            "Edge sparsification — permanently delete a fraction of edges once, before training.",
        ],
        notes=(
            "These are practical knobs teams might try before DP training. Stress that hyperparameters are fixed across "
            "the grid for comparability. DropEdge is stochastic per epoch; sparsification is a one-time graph edit. "
            "Masking mimics APIs that only return top classes. Foreshadow that label smoothing and sparsification behave "
            "non-monotonically for membership risk."
        ),
    )

    section_slide(
        prs,
        "Results",
        "Synthetic graphs: where structure shows up",
        notes=(
            "Transition slide: synthetic data is not “more real,” but it is where you can isolate homophily and density. "
            "Tell the audience to expect near-chance AUROC under high homophily and a spike for GCN under low homophily "
            "and sparsity, with GraphSAGE often remaining closer to chance despite similar accuracy in many cells."
        ),
    )

    takeaway_slide(
        prs,
        "When neighbors mostly disagree with the label (low homophily), GCN outputs can become easier to separate for members vs non-members—especially when the graph is sparse.",
        notes=(
            "Interpretation hook: heterophily makes neighbor mixing a noisy operation; GCN’s smoothing can create "
            "systematic confidence differences between supervised and held-out nodes that a linear attacker can exploit. "
            "Sparse graphs amplify the effect because each node depends on fewer, more pivotal edges."
        ),
        supporting="Attack AUROC near 0.5 means “about as hard as guessing a coin flip.” Higher means the attack is working.",
    )

    figure_slide(
        prs,
        "No defense: who looks most “memorized”?",
        "basic_fig1_attack_auc_models_no_defense.png",
        "Low homophily + sparse edges: GCN stands out (~0.59 AUROC). High homophily: everyone sits near random chance (~0.5).",
        notes=(
            "Walk left-to-right across regimes on the x-axis. Contrast the tall GCN bar in low-homophily sparse settings "
            "with the cluster near 0.5 in high homophily. Point out LogReg/MLP/GraphSAGE remain comparatively flat—this "
            "is the core structural story before introducing defenses. Mention error bars reflect five seeds."
        ),
        caption="Bars = mean attack AUROC across seeds; error bars show spread.",
    )

    figure_slide(
        prs,
        "Same story, one heatmap",
        "basic_fig4_heatmap_models_x_setting.png",
        "The pattern is not “GNNs are always worse”—it is regime-specific.",
        notes=(
            "Use this as a compact visual summary. Rows are model families; columns are synthetic regimes. Ask the "
            "audience to notice that leakage is localized to certain cells rather than uniformly high for all GNNs. "
            "This supports the ‘no monolithic GNN risk’ message."
        ),
        caption="Rows are models; columns are synthetic regimes.",
    )

    figure_slide(
        prs,
        "Zoom in: GCN only",
        "basic_fig2_gcn_attack_auc_by_setting.png",
        "For GCN without defense, leakage rises when homophily drops; within low homophily, sparser graphs often leak more in this benchmark.",
        notes=(
            "Isolate GCN so the professor sees the homophily and density trend without comparing across models. "
            "If asked about mechanism, tie back to coupled optimization and how heterophilous neighborhoods change "
            "train vs held-out posteriors under spectral-style aggregation."
        ),
        caption=None,
    )

    two_caption_slide(
        prs,
        "Utility vs privacy (and homophily)",
        "Good node accuracy does not automatically mean safe posteriors—defenses can move calibration and leakage in different directions.",
        (
            "fig2_privacy_utility_tradeoff.png",
            "Each point: one (model, defense) run—accuracy vs attack AUROC.",
        ),
        (
            "fig3_leakage_vs_homophily.png",
            "Measured homophily vs attack AUROC across synthetic runs.",
        ),
        notes=(
            "Left plot: each point is a defense/model combination—show that moving right (higher attack AUROC) does not "
            "always align with lower accuracy. Right plot: relate realized homophily on the x-axis to leakage on the y-axis. "
            "This pair supports the ‘utility parity is not privacy parity’ discussion, especially for GraphSAGE vs GCN."
        ),
    )

    figure_slide(
        prs,
        "Density matters too",
        "fig5_leakage_vs_density.png",
        "After fixing homophily, edge budget (density) still shifts how much signal the attack sees.",
        notes=(
            "Explain ρ as the fraction of possible edges realized before sampling rejects duplicates. Within a homophily "
            "band, changing density changes receptive field size and can change how sharply members separate from "
            "non-members in softmax space. Keep claims empirical—this is an observed trend in your grid, not a theorem."
        ),
        caption=None,
    )

    section_slide(
        prs,
        "Citation graphs",
        "When features steal the spotlight",
        notes=(
            "Pivot slide: moving from controlled synthetics to Planetoid benchmarks where bag-of-words features are highly "
            "informative. Prepare the audience for a reversal—feature-only models can exhibit higher attack AUROC than "
            "GCN even though the story on synthetics emphasized GCN leakage."
        ),
    )

    takeaway_slide(
        prs,
        "On Cora and Citeseer, LogReg and MLP can leak more than GCN—not because “graphs are safer,” but because easy-to-separate features can dominate the attack signal.",
        notes=(
            "Emphasize confounding: citation graphs mix strong linear separability with graph structure. A GNN’s mixing "
            "can sometimes smooth confidence patterns that a linear model exaggerates. This is not a claim that graphs "
            "inherently protect privacy—only that the dominant leakage driver is setting-dependent."
        ),
    )

    figure_slide(
        prs,
        "Cora & Citeseer: no defense",
        "fig1_attack_auc_vs_model.png",
        "Example numbers your audience can remember: strong linear models on bag-of-words features can look “more memorized” than a GCN in softmax space.",
        notes=(
            "Point to concrete bars: on Cora, MLP and LogReg exceed GCN for the confidence attack in your reported run. "
            "On Citeseer, MLP is especially high. Use this to motivate measuring membership risk on the actual deployment "
            "graph and model rather than assuming GNNs are always worse."
        ),
        caption="Compare bars within each dataset; read vertically for model type.",
    )

    section_slide(
        prs,
        "Defenses",
        "Why “obvious” fixes can backfire",
        notes=(
            "Set expectations: none of these defenses are privacy guarantees. You will highlight label smoothing as "
            "often increasing attack AUROC, and sparsification as not providing a robust win once the same thinned graph "
            "is used consistently for training, inference, and shadow evaluation."
        ),
    )

    takeaway_slide(
        prs,
        "Label smoothing often made membership attacks easier in our experiments—not harder. Do not treat it as a privacy knob without measuring attacks.",
        notes=(
            "Explain the intuition gap: smoothing can reshape posteriors without making members and non-members more "
            "alike under graph coupling. Cite paired-test examples if asked (GraphSAGE on low-hom sparse synthetics, "
            "Cora GCN, Citeseer GraphSAGE). Encourage reporting ECE alongside AUROC because calibration can move "
            "independently from membership separability."
        ),
        supporting="Edge sparsification did not produce a reliable privacy win for GCN once we used the same thinned graph for training, testing, and shadow models.",
    )

    figure_slide(
        prs,
        "Example: GCN + defenses (one regime)",
        "basic_fig3_gcn_defenses_low_med.png",
        "Defenses move the needle differently by mechanism—there is no single winner on every bar.",
        notes=(
            "Walk through each colored/defense bar for GCN on low-homophily medium-density synthetics. Note some defenses "
            "help a little, some hurt, and none universally collapse AUROC to chance. This is the ‘evaluate per regime’ "
            "message in one panel."
        ),
        caption="Low homophily, medium density synthetic regime.",
    )

    figure_slide(
        prs,
        "Citations: defense ranking shifts",
        "fig4_defense_effectiveness.png",
        "The “best” defense for privacy depends on the dataset and model—rankings are not portable by eyeball alone.",
        notes=(
            "Use this slide if the professor asks about real-world relevance. Rankings differ by dataset and architecture; "
            "avoid claiming one defense is always best. Offer to open the comprehensive heatmap next for detailed Q&A."
        ),
        caption=None,
    )

    figure_slide(
        prs,
        "Full experimental grid",
        "fig6_comprehensive_heatmap.png",
        "Use this as the “coverage map” for Q&A: many cells, one consistent pipeline.",
        notes=(
            "This is your ‘coverage’ slide during Q&A. It communicates breadth: many dataset×model×defense combinations "
            "with one protocol. If pressed on multiple comparisons, remind them p-values were exploratory and more seeds "
            "would strengthen statistical claims."
        ),
        caption=None,
    )

    bullet_slide(
        prs,
        "Remember these three points",
        [
            "Privacy risk is contextual: it comes from the interaction of graph structure, model architecture, features, and the defense—not from the word “GNN” alone.",
            "High accuracy does not imply low membership risk; two models can match on accuracy but differ a lot in attack AUROC.",
            "Measure membership and calibration together; a defense can change probabilities in ways that help attacks even when accuracy looks fine.",
        ],
        notes=(
            "Close the scientific story before limitations. Offer one sentence per bullet aloud, slowly. This is what you "
            "want remembered after the talk ends. Optionally connect to responsible deployment: define the membership "
            "predicate, log confidence histograms, and run attack evaluations when changing models or defenses."
        ),
    )

    bullet_slide(
        prs,
        "Honest limits of this study",
        [
            "Synthetic graphs simplify the world (e.g., no rich community structure).",
            "We only evaluate score-based black-box attacks—stronger audits are future work.",
            "Five seeds means statistics are illustrative; many comparisons were run, so treat p-values as exploratory.",
        ],
        notes=(
            "Anticipate methodology questions. Acknowledge generator simplicity and the lack of LiRA-style or low-FPR "
            "audits. Mention DP-GNN literature as the rigorous alternative when regulatory guarantees are required. "
            "Five seeds is enough for a student-scale study but not for definitive inference—frame statistics honestly."
        ),
    )

    bullet_slide(
        prs,
        "Thank you",
        [
            "Repository: privacy-gnn (YAML configs, PyTorch Geometric, CSV outputs, figure scripts).",
            "Happy to dive into any cell of the heatmaps or the attack feature vector φ.",
        ],
        notes=(
            "Offer concrete follow-ups: walk through a single CSV row, show how φ is computed from logits, or demo "
            "re-running one YAML configuration. Thank your professor by name if appropriate when presenting live."
        ),
    )

    prs.save(OUT)
    print(f"Wrote: {OUT}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
